# app/parsers.py
import os, re
from typing import List, Tuple

from markdown_it import MarkdownIt
md = MarkdownIt()

def read_markdown_text(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def read_pdf_text(path: str) -> str:
    import fitz  # PyMuPDF
    doc = fitz.open(path)
    parts = []
    for p in doc:
        t = p.get_text("text")
        if t:
            parts.append(t.strip())
    doc.close()
    return "\n\n".join(parts)

def read_pptx_text(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    chunks = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        if slide.shapes.title and slide.shapes.title.text:
            texts.append(slide.shapes.title.text)
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        slide_text = "\n".join([t.strip() for t in texts if t and t.strip()])
        if slide_text:
            chunks.append(f"# Slide {i}\n{slide_text}")
    return "\n\n".join(chunks)

def read_docs(paths: List[str]) -> List[Tuple[str, str]]:
    """
    Returns list of (basename, full_text)
    Supports .md, .markdown, .pdf, .pptx
    """
    out = []
    for p in paths:
        ext = os.path.splitext(p)[1].lower()
        base = os.path.basename(p)
        text = ""
        try:
            if ext in (".md", ".markdown"):
                text = read_markdown_text(p)
            elif ext == ".pdf":
                text = read_pdf_text(p)
            elif ext == ".pptx":
                text = read_pptx_text(p)
            else:
                # unknown: best effort read as text
                text = read_markdown_text(p)
        except Exception as e:
            text = f""  # skip on failure
        if text and text.strip():
            out.append((base, text))
    return out
